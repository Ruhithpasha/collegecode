from pptx import Presentation

prs = Presentation()

slides = [
    ("Nirvanix – Cloud Storage Platform", "Overview, Offerings, Pros & Cons\nPresented by: [Your Name]"),
    ("What is Nirvanix?", "• Cloud storage company for enterprise-level needs\n• Provided global file access via distributed data centers\n• Known for Storage Delivery Network (SDN)"),
    ("Why Should We Use Nirvanix?", "• Enterprise-grade and scalable\n• Strong encryption and compliance\n• Global data access with low latency"),
    ("Offerings of Nirvanix", "• CloudNAS – mount as local drive\n• SDN – Storage Delivery Network\n• Hybrid cloud integration (hNode)\n• APIs and SDKs\n• Backup and archiving solutions"),
    ("When It Started", "• Founded in 2007 in San Diego, California\n• Evolved from StreamLoad\n• Raised over $70 million in funding"),
    ("Why Did Nirvanix Fall?", "• Leadership instability\n• High operational costs\n• Strong competition from AWS, Azure\n• Limited time for customer data migration\n• Filed for bankruptcy in October 2013"),
    ("Advantages of Nirvanix", "• Enterprise-level features\n• Global access and scalability\n• Strong security\n• Hybrid cloud flexibility"),
    ("Disadvantages of Nirvanix", "• High cost\n• Low brand recognition\n• Sudden shutdown risk\n• Less developer-friendly than competitors"),
    ("Who Could Use Nirvanix?", "• Large enterprises with big data\n• Media and entertainment companies\n• Backup and disaster recovery services\n• Developers of data-intensive applications"),
    ("Conclusion", "• Early innovator in cloud storage\n• Fell due to financial and competitive issues\n• Important lessons in cloud service management")
]

for title, content in slides:
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = content

prs.save("Nirvanix_Presentation.pptx")
